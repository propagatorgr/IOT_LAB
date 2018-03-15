##### What's News App ####
# http://python-pptx.readthedocs.io/en/latest/user/quickstart.html
# sudo pip3 install feedparser
# sudo apt-get install python3-lxml
# http://python-pptx.readthedocs.io/en/latest/user/install.html
# https://www.raspberrypi.org/forums/viewtopic.php?t=86441

import feedparser
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

### email modules ###
import smtplib,ssl
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email.mime.text import MIMEText
from email.utils import formatdate
from email import encoders

### lists to hold headline and link ###
News = []
Links = []

### Function to send the email ###
def send_an_email():
    toaddr = '########' # Η διεύθυνση του παραλήπτη   
    me = '############' # Η διεύθυνση του αποστολέα
    subject = "What's News"

    msg = MIMEMultipart()
    msg['Subject'] = subject
    msg['From'] = me
    msg['To'] = toaddr
    msg.preamble = "test " 
    #msg.attach(MIMEText(text))

    part = MIMEBase('application', "octet-stream")
    part.set_payload(open("Whats_News.pptx", "rb").read())
    encoders.encode_base64(part)
    part.add_header('Content-Disposition', 'attachment; filename="Whats_News.pptx"')
    msg.attach(part)

    try:
       s = smtplib.SMTP('smtp αποστολέα', smtp number)
       s.ehlo()
       s.starttls()
       s.ehlo()
       s.login(user = 'Η διεύθυνση αποστολέα', password = 'Συνθηματικό αποστολέα')
       #s.send_message(msg)
       s.sendmail(me, toaddr, msg.as_string())
       s.quit()
    except:
       print ("Error: unable to send email")
    #except SMTPException as error:
         # print ("Error")

### Create News Feed, pull down latest News ###
Esosnews = feedparser.parse("https://www.esos.gr/rss.xml")
#"https://www.esos.gr/rss.xml"
### Create the slide layout ###
def create_my_default_slide(title,subtitle):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    
    return slide

### create the hyperlink layout ###
def add_hyper_link(shape, text, url):
    p = shape.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    p.font.size = Pt(12)
    p.font.name = 'Verdana'
    p.font.color.rgb = RGBColor(0, 0, 255)
    hlink = r.hyperlink
    hlink.address = url

### pull down the news stories ###
for i in range(10):
    text = Esosnews["entries"][i]["title"]
    desc = Esosnews["entries"][i]["description"]
    links = Esosnews["entries"][i]["link"]
    print(text)
    print(links)
    print(desc)
    News.append(text) #add headlines to News list
    Links.append(links) # add links to Links list

# print (News)
# print (Links)

### Create the Dictionary from the lists News and Links ###
News_Dict = {}
for i in range(len(News)):
    News_Dict [News[i]] = Links[i]
print ("")    
print (News_Dict)    #test#

###### Create a Powerpoint ######
### Loops over keys and values  in the Dictionary and adds them to the slide ###
### Set up PPT ###
prs = Presentation()
for key, value in News_Dict.items():
    this_slide = create_my_default_slide("%s" % key," \n  \n  \n ")
    add_hyper_link(this_slide.shapes[1], value, value)

print(len(prs.slides))

### save the Powerpoint ###
prs.save('Whats_News.pptx')
send_an_email()
